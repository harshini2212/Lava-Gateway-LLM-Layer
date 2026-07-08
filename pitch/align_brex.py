"""Slight, truthful tweaks to slides 2, 6, 11 so the language maps onto Brex's core —
spend management, cash flow, burn, runway — without overstating the candidate's real work."""
import io, shutil
from pptx import Presentation

SRC = r"C:/Users/harsh/Downloads/Brexify-Jobs-Brex.pptx"
OUT = "brexify-pitch-final.pptx"

EDITS = [
    # Slide 2 — thesis + the two work cards
    ("Google Shape;67;p10", 0, 0, "making AI output over spend & cash flow "),
    ("Google Shape;74;p10", 1, 0, "Causal cash-flow & runway engine for CFOs"),
    ("Google Shape;85;p10", 1, 0, "Agentic causal-graph engine for burn & cash-runway"),
    # Slide 6 — founder description + MVP milestone
    ("Google Shape;174;p14", 0, 0,
     "Causal cash-flow & runway simulation engine for CFOs. Completed customer discovery calls and "),
    ("Google Shape;187;p14", 1, 1, " Architected the core cash-flow engine."),
    # Slide 11 — closing tagline (same length as original, zero overflow risk)
    ("Google Shape;256;p19", 0, 0, "Let's make spend & cash flow AI actually trustworthy."),
]


def index_shapes(shapes, out):
    for sh in shapes:
        out[sh.name] = sh
        if sh.shape_type == 6:
            index_shapes(sh.shapes, out)
    return out


prs = Presentation(SRC)
names = {}
for slide in prs.slides:
    index_shapes(slide.shapes, names)

applied, missing = 0, []
for name, p, r, text in EDITS:
    sh = names.get(name)
    if sh is None or not sh.has_text_frame or p >= len(sh.text_frame.paragraphs) \
            or r >= len(sh.text_frame.paragraphs[p].runs):
        missing.append((name, p, r))
        continue
    sh.text_frame.paragraphs[p].runs[r].text = text
    applied += 1

prs.save(OUT)
shutil.copyfile(OUT, SRC)  # overwrite the Downloads deliverable
print(f"applied {applied}/{len(EDITS)} edits; missing {missing}")
