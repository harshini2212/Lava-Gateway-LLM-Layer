"""Rewrite the Hebbia job-pitch deck into a Lava / Lavagent pitch.

Only company-specific content changes (slides 1,2,3,4,11). Personal slides (5-10:
experience, founder story, hackathons, GitHub) are untouched. Edits are applied at the
run level so every font, color, size, position and image is preserved.
"""
import io
from pptx import Presentation

SRC = "hebbia-original.pptx"
OUT = "lavagent-pitch.pptx"

# (shape_name, paragraph_index, run_index, new_text). '' blanks a run.
EDITS = [
    # ---- Slide 1: title / live link ----
    ("Google Shape;50;p9", 1, 0, "lavagent.up.railway.app"),

    # ---- Slide 2: Why Lava? + roles ----
    ("Google Shape;64;p10", 0, 0, "Why Lava?"),
    ("Google Shape;66;p10", 0, 0,
     "The reason Lava caught my attention: a lot of my recent work is about"),
    ("Google Shape;88;p10", 1, 0, "Software "),
    ("Google Shape;88;p10", 1, 1, "Engineer, AI & Agents"),
    ("Google Shape;88;p10", 1, 2, ""),
    ("Google Shape;88;p10", 1, 3, ""),
    ("Google Shape;88;p10", 2, 0, "Forward Deployed Engineer"),
    ("Google Shape;88;p10", 2, 1, ""),

    # ---- Slide 3: Case Study: Lavagent ----
    ("Google Shape;95;p11", 0, 0, "Case Study: Lavagent "),
    ("Google Shape;97;p11", 1, 1, "Lavagent"),
    ("Google Shape;97;p11", 1, 2,
     ", a financial-correctness harness for Lava spend that scores every transaction "
     "for fraud, replays policy, and tie-out-verifies every number"),
    ("Google Shape;97;p11", 1, 3,
     " — so Lava can scale autonomous finance agents without the risk of approving a "
     "fraudulent charge."),
    ("Google Shape;97;p11", 3, 2, "lavagent.up.railway.app  "),
    ("Google Shape;100;p11", 0, 0,
     "By enforcing spend policy and accounting identities at the decision layer, the "
     "harness catches fraud, duplicate invoices, and policy violations before money "
     "moves — the integrity Lava needs to let agents act on real dollars."),
    ("Google Shape;105;p11", 0, 0, "Policy + Identities"),
    ("Google Shape;105;p11", 1, 0, "spend within limit · books tie out"),
    ("Google Shape;105;p11", 2, 0, "✔ Enforced"),
    ("Google Shape;108;p11", 1, 0, "Three-way match · duplicate & bank-change detection"),
    ("Google Shape;111;p11", 1, 0, "Calibrated Lava tenant · fraud ROC-AUC 0.89+"),

    # ---- Slide 4: Card + Cash + ERP ----
    ("Google Shape;117;p12", 0, 0, "Card + Cash + ERP: exactly how they connect"),
    ("Google Shape;119;p12", 0, 0,
     "Lava Assistant’s real job is unifying card spend, cash, and the GL into one "
     "ledger — checked by the same engine."),
    ("Google Shape;121;p12", 0, 0, "\U0001F4C4  Books"),
    ("Google Shape;121;p12", 1, 0, "GL · ERP (NetSuite)"),
    ("Google Shape;123;p12", 0, 0, "\U0001F5C4  Lava"),
    ("Google Shape;123;p12", 1, 0, "Card · Lava Cash · AP"),
    ("Google Shape;126;p12", 0, 0, "Lavagent"),
    ("Google Shape;126;p12", 1, 0, "scores every txn · replays policy"),
    ("Google Shape;126;p12", 2, 0, "reconcile to the GL?"),
    ("Google Shape;129;p12", 1, 0, "close-ready & audit-ready"),
    ("Google Shape;130;p12", 0, 0,
     "A synthetic Lava tenant is calibrated to realistic spend — fraud, policy, and "
     "tie-out all proven on it. “Calibrated demo data.”"),
    ("Google Shape;132;p12", 0, 0, "Card + AP unlocks"),
    ("Google Shape;133;p12", 0, 1, "Fraud rings "),
    ("Google Shape;133;p12", 0, 2, "— shared-card collusion the static rules miss."),
    ("Google Shape;133;p12", 1, 1, "Policy leakage "),
    ("Google Shape;133;p12", 1, 2, "— replay a new policy over history, see the $ impact."),
    ("Google Shape;133;p12", 2, 1, "Duplicate & bank-change AP "),
    ("Google Shape;133;p12", 2, 2, "— stopped before Bill Pay sends a cent."),
    ("Google Shape;135;p12", 0, 0, "Cash + Credit unlocks"),
    ("Google Shape;136;p12", 0, 1, "Causal runway "),
    ("Google Shape;136;p12", 0, 2, "— do-operator what-ifs, not a re-plotted trend."),
    ("Google Shape;136;p12", 1, 1, "Idle-cash yield "),
    ("Google Shape;136;p12", 1, 2, "— the Lava Cash sweep opportunity, quantified."),
    ("Google Shape;136;p12", 2, 1, "Underwriting & limits "),
    ("Google Shape;136;p12", 2, 2, "— a PD model recommends the credit line."),
    ("Google Shape;138;p12", 0, 0,
     "The same engine that scores a charge proves the books tie out to it — "),
    ("Google Shape;138;p12", 0, 1,
     "the trust layer that lets Lava Assistant act on real spend."),

    # ---- Slide 11: Why Lava closing ----
    ("Google Shape;250;p19", 0, 1, "LAVA"),
    ("Google Shape;251;p19", 0, 2, "Lava "),
]


def index_shapes(shapes, out):
    for sh in shapes:
        out[sh.name] = sh
        if sh.shape_type == 6:  # group
            index_shapes(sh.shapes, out)
    return out


prs = Presentation(SRC)
name_map = {}
for slide in prs.slides:
    index_shapes(slide.shapes, name_map)

applied, missing = 0, []
for shape_name, p, r, text in EDITS:
    sh = name_map.get(shape_name)
    if sh is None or not sh.has_text_frame:
        missing.append((shape_name, p, r, "no shape"))
        continue
    paras = sh.text_frame.paragraphs
    if p >= len(paras) or r >= len(paras[p].runs):
        missing.append((shape_name, p, r, f"p/r out of range (paras={len(paras)})"))
        continue
    paras[p].runs[r].text = text
    applied += 1

prs.save(OUT)
log = io.open("edit_log.txt", "w", encoding="utf-8")
log.write(f"applied {applied}/{len(EDITS)} edits -> {OUT}\n")
for m in missing:
    log.write("MISSING: " + repr(m) + "\n")
log.close()
print(f"applied {applied}/{len(EDITS)}; missing {len(missing)}")
